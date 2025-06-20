import os
import subprocess
from django.conf import settings
from django.http.response import StreamingHttpResponse
from django.shortcuts import redirect, render
from django.http import FileResponse, Http404, JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from .utils import generate_response
from django.contrib.auth import get_user_model
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import re
import io
import uuid
import threading
from django.contrib.auth.decorators import login_required
from django.views.decorators.cache import cache_control


# def indexs(request):
#     if 'user' in request.session:
#         db = get_user_model()
#         current_user = request.session['user']
#         isUserAdmin = db.objects.filter(username = current_user, is_superuser = 1).exists()
#         params = {'current_user': current_user, "isAdmin": isUserAdmin}
#         return render(request, 'kunwariwebpage/index.html', {"params":params})
#     else:
#         return redirect('login')
    
@login_required
@cache_control(no_cache=True, must_revalidate=True, no_store=True)
def index(request):
    try:
        current_user = request.user
        isUserAdmin = current_user.is_superuser  # no need to query DB again

        params = {
            'current_user': current_user.username,
            'isAdmin': isUserAdmin
        }
        return render(request, 'kunwariwebpage/index.html', {"params": params})
    except:
        return JsonResponse({"error": "Login failed"}, status=500)
    


#def register(request):
#    return render(request, 'kunwariwebpage/registration.html')

#def forgotpassword(request):
#    return render(request, 'kunwariwebpage/forgotpass.html')

def emu_to_inches(emu):
    return emu / 914400

@csrf_exempt
def chatbot_view(request):
        try:
            if request.method == "POST":
                message = request.POST.get("message")
                files = request.FILES.getlist("files")
                print(message)
                print(files)
                #data = json.loads(request.body)
                #print(data)
                #message = data['message']
                message += """
### Lessons - Set lesson in a storyworld where learning activities become quests, and students take on roles, and the teacher becomes the main storyteller or quest master.

---

### Lectures - Automatically Create short lectures thought through analogies consistent with the story world

---

### Teaching and learning activities - Automatically Create Teaching and learning activities aligned with the story world and the quests actively engages learners through multisensorial learning

---

### Learning Objectives (Bloom Verb-Based) to Verb-Based Game Mechanics - Automatically transform the Learning Objectives into Game Mechanics that integrates with the plot of the key story. 
##Objective 1: Learning Objective 1 tied to episode 2 of the story stated in a verb-based quest format
##Objective 2: Learning Objective 2 tied to episode 2 of the story stated in a verb-based quest format
##Objective 3: Learning Objective 3 tied to episode 3 of the story stated in a verb-based quest format, which serves as the Grand Quest

---

### Assessments - Automatically create gamified assessments  aligned with the learning objectives integrated within the story.

---

### Key moral lesson/values(Narrative) - Automatically generate the Morals or Values arc connected to the key story/premise.

---

### Player Types (Segmentation). All the specific activities should be part of the story.
Killer: Automatically Create Specific activities for competitive players.  Achiever: Automatically Create Specific activities Focused on accomplishment-based challenges. Explorer: Automatically creates specific activities Engages with discovery-based elements. Socializer: Automatically Create Specific Activities Collaborative and team-oriented tasks. 
Life Skills, Soft Skills, Creative Skills, 5 C’s of 21st-century learning Application - Automatically tailor the games according to applicable skills. 

---

### Key Resources - Automatically generate the following: Props, Manipulatives, and Learning Materials: Suggestions for physical or virtual game materials. Background/Set: Virtual or physical space design ideas. Costume/Attire: Suggestions for character costumes or thematic attire. 

---

### Activities - Automatically generate the following: Dance/Music/SFX: Tailored sound design suggestions for immersion. Food, Taste, and Scents: Olfactory enhancements for deeper engagement.

---

### Student/Teacher Roles (Kalaro)  - Assigns roles to students, teachers, and NPCs in alignment with the key story.

---

### Reflection and Discussion - Discussion Questions: Auto-generated to facilitate meaningful reflection. 

---

### Main Challenge (Summative Assessment)  - Create a final stage with a BOSS challenge that integrates all learning outcomes.

---

### What’s In It For You?/ Why is this relevant? - Automatically generate answers to this question from real-world scenarios or UN SDG or the 5Cs of 21st century Learning

---

### Bonus Challenge - Automatically generate a challenge from real-world scenarios or UN SDG that students can solve from what they’ve learned from the topic - MAKE THIS PISA CREATIVE THINKING FORMAT 

---

### Rewards and Badges - Automatically unlock Achievements based on performance and engagement.

---
"""
                # with open("static/Kunwari Activator Template (KATE) V2.pdf", "r", encoding="utf-8") as f:
                #      instructions = f.read()
                # message+= instructions
                response = StreamingHttpResponse(generate_response(message, files), status=200, content_type='text/plain')
                #return render(request, "kunwariwebpage/index.html", {"response": 'response'})
                return response
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
        return JsonResponse({"error": "Invalid request"}, status=400)

@csrf_exempt
def generate_episodes(request):
     try:
        if request.method == "POST":
            data = json.loads(request.body)
            #Message should be the response of the bot + the Episodes
            message = data['message']
            message += """
Follow this rules when generating the episode:
*Maintain narrative continuity with the key story/premise
*Specify time allocations and technology requirements
*Suggest physical and digital resource alternatives

### Episode 1: "Concept Introduction" - Create interactive slides that:
*Introduce the key concepts and foundational knowledge
*Establish the narrative framework and character roles
*Include 3-5 concrete examples that contextualize the content
*End with a mini-challenge that tests basic concept understanding

### Episode 2: "Skill-Building Activities" - Develop activity-focused slides that:
*Present 2-3 scaffolded activities for each key concept
*Incorporate player-type variations (competitive, achievement, exploratory, collaborative)
*Connect activities directly to the story and to real-world applications

### Episode 3: "Applied Learning Assessments" - Design assessment slides that:
*Frame assessments as narrative missions with clear objectives
*Include both individual and group assessment options
*Provide rubrics that align with learning objectives
*Offer differentiated assessment paths based on student strengths

### Episode 4: "Critical Reflection & Extension" - Create culminating slides that:
*Guide structured reflection on content application
*Present the "boss challenge" with real-world relevance
*Incorporate societal impact connections and 21st century skills application
*Include extension resources for continued exploration
"""
            
            """Please transform the lesson plan above following the format of Episodes:
Episode 1: Gamified Lessons - revise the slides from lesson plan.
Episode 2: Gamified Teaching and Learning Activities - Create slides for interactive activities. 
Episode 3: Gamified Assessments - Create slides for narrative-based assessments. 
Repeat as needed for further episodes -  Tailored lessons, activities, and assessments in sequence.
            """

            """Please transform the lesson plan above following the format of Episodes:
Episode 1: Gamified Lessons - automatically create slides from lesson modules.
Episode 2: Gamified Teaching and Learning Activities - Create slides for interactive activities. 
Episode 3: Gamified Assessments - Create slides for narrative-based assessments. 
Repeat as needed for further episodes -  Tailored lessons, activities, and assessments in sequence.

Strictly follow the time duration mentioned
            """
            response = StreamingHttpResponse(generate_response(message), status=200, content_type='text/plain')
            return response
     except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
     return JsonResponse({"error": "Invalid request"}, status=400)

@csrf_exempt
def generate_content(request):
     try:
        if request.method == "POST":
            data = json.loads(request.body)
            message = data['message']
            message += """
Please generate comprehensive slide content for each episode, including:


## Episode 1: "Concept Introduction" - Create detailed slides with(separate each slide with ---):
Opening slide with compelling hook and visual theme
Content slides for each key concept (5-8 slides per concept)
Visual examples and case studies with explanatory text
Interactive elements and discussion prompts
Mini-challenge slide with instructions and assessment criteria

---

## Episode 2: "Skill-Building Activities" - Develop activity slides with(separate each slide with ---):
Activity overview and learning connection for each activity
Step-by-step instructions with visuals
Player-type variations with differentiated instructions
Printable worksheets and digital templates
Activity debriefing questions

---

## Episode 3: "Applied Learning Assessments" - Design assessment slides with(separate each slide with ---):
Assessment narrative framing and mission objectives
Clear assessment criteria and expectations
Individual and group assessment options
Rubric slides with performance indicators
Self-assessment components

---

## Episode 4: "Critical Reflection & Extension" - Create reflection slides with(separate each slide with ---):
Guided reflection prompts and documentation templates
Boss challenge instructions and success criteria
Real-world application connections
Extension options and resources
Celebration and accomplishment recognition

---
"""  
            """
            (Please include images(Online URLs) for each topic that are accessible online Following the format ![alt text](Image URL) and make sure it was on a new line)
            Please include images(Online URLs) for each topic that are accessible online based on the headings. Following the format ![alt text](Image URL) and make sure it was on a new line,
            Please generate detailed content for the Episode slides above."""
            response = StreamingHttpResponse(generate_response(message), status=200, content_type='text/plain')
            return response
     except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
     return JsonResponse({"error": "Invalid request"}, status=400)

@csrf_exempt
def generate_facilitator_script(request):
     try:
        if request.method == "POST":
            data = json.loads(request.body)
            message = data['message']
            message += """
Based on the episodic lesson plan developed, create comprehensive facilitator script for effective implementation:

Facilitator Overview:
  Provide a 1-page summary of the entire lesson sequence
  Include learning objectives, and required preparation
  Outline technology and resource needs with alternatives

---

For Each Episode (1-4), Create Detailed Scripts.
Format all materials for easy reference during implementation, with clear section headings, visual cues, and a consistent layout that distinguishes between:
*Essential instructions (must do)
*Optional enhancements (could do)
*Troubleshooting tips (if needed)

---

Based on the episodic lesson plan developed, create comprehensive facilitator materials for effective implementation:

### Facilitator Overview
*Provide a 1-page summary of the entire lesson sequence
*Include learning objectives, time requirements, and required preparation
*Outline technology and resource needs with alternatives

---

For Each Episode (1-4), Create Detailed Implementation Guides:
### Pre-Session Preparation
- Materials checklist (digital and physical)
- Room/space setup instructions
- Technology setup and troubleshooting tips
- Pre-session facilitator knowledge check

---

### Step-by-Step Facilitation Script
- Opening hooks and engagement strategies
- Transition scripts between activities
- Sample dialogue with anticipated student responses
- Time allocations for each segment (minimum/ideal)
- Adaptation notes for different learning contexts

---

### Student Support Guide
- Common misconceptions and how to address them
- Differentiation strategies for various learning needs
- Extension prompts for advanced learners
- Support scaffolds for struggling learners

---

### Assessment Implementation
- Detailed evaluation criteria and success indicators
- Observation prompts during activities
- Feedback delivery scripts
- Documentation strategies

---

### Reflection and Closure
- Guided discussion questions with sample responses
- Consolidation strategies for key concepts
- Bridge statements to subsequent episodes
- Follow-up activities and homework options

---
"""
            """Please generate a complete facilitator script for the Episode slides above."""
            response = StreamingHttpResponse(generate_response(message), status=200, content_type='text/plain')
            return response
     except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
     return JsonResponse({"error": "Invalid request"}, status=400)


def generate_pptx(request):
    prs = Presentation()
    pptx_io = io.BytesIO()
    # dimension_width = emu_to_inches(prs.slide_width)
    # dimension_height = emu_to_inches(prs.slide_height)
    # left_offset = 1
    # top_offset = 2
    try:
        if request.method == "POST":
            data = json.loads(request.body)
            message = data['message']
            fileName = data['filename']
            parts = re.split('<hr>', message)
            """
            Find the HTML tags =  r'<(h[1-3]|p)>(.*?)</\1>', ['h4', 'Key Plot/Conflict/Obstacle']
            Use the index[0] as indicator
                If tag is h1,h2,h3 use layout[5] Title only
                If tag is h4,h5,h6 use layout[1] Title and content
                If tag is p, insert it to layout[1]
            Finally, trim each sentence before inserting to pptx
            """
            pattern = r'<([a-zA-Z][a-zA-Z0-9]*)[^>]*>(.*?)</\1>'
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            content_placeholder = ""
            run_once = True
            for part in parts:
                if part == " ":
                    print("Part is None")
                matches = re.findall(pattern, part) #[('h1', 'Title One'), ('p', 'This is a paragraph.'), ('h2', 'Subtitle'), ('h3', 'Section')]
                if run_once == False and len(matches) != 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    content_placeholder = slide.placeholders[1]
                for i in matches: 
                    if run_once:
                        title_slide.shapes.title.text = i[1]
                        run_once = False
                    #['h4', 'Key Plot/Conflict/Obstacle']
                    elif i[0] in ['h1','h2', 'h3', 'h4', 'h5', 'h6']:
                        #print('Title Match')
                        slide.shapes.title.text = i[1]
                    else:
                        trim_part = re.sub(r'<[^>]+>', '', i[1])
                        content_placeholder.text += f"{trim_part}\n"
                        
            #use file-like object instead of saving in path
            #prs.save(os.path.join(settings.MEDIA_ROOT,fileName))
            #return JsonResponse({"message": "File generated", "filename": fileName}, status=200)

            # Use BytesIO to save the presentation in memory
            prs.save(pptx_io)
            pptx_io.seek(0)  # Important: go to the beginning of the stream

            # Create response
            # response = HttpResponse(content_type='application/vnd.ms-powerpoint')
            # response['Content-Disposition'] = 'attachment; filename="sample.pptx"'
            # response.write(pptx_io.getvalue())
            # pptx_io.close()
            # return response
            response = HttpResponse(
                pptx_io.read(),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = 'attachment; filename="sample.pptx"'
            #print(response)
            return response
    except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    
def generate_pptxs(request):
    try:
        if request.method != "POST":
            return  JsonResponse({"error": "Invalid request"}, status=405)
        
        #title = request.POST.get('title')
        #content = request.POST.get('content')
        data = json.loads(request.body)
        message = data['message']
        fileName = str(uuid.uuid4())
        #fileName = data['filename']

        qmd_content = f"""---\ntitle: "Sample powerpoint"\nformat: pptx\nauthor: "ivan"\n---\n{message}"""
        trim_tabs = re.sub(r'[ \t]*\n[ \t]*', '\n', qmd_content)
        cleaned_text = re.sub(r'^[#\*]+\s*slide \d+:\s*', '', trim_tabs, flags=re.IGNORECASE | re.MULTILINE)
        #print(repr(cleaned))
        # clean up the qmd file, remove extra tabs(\t) for each new line
        input_qmd_path = os.path.join(settings.BASE_DIR, "quarto_output", "files" ,f'{fileName}.qmd')
        output_pptx_path = os.path.join(settings.BASE_DIR,"quarto_output", "files", f'{fileName}.pptx')

        #os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
        
        with open (input_qmd_path, 'w', encoding="utf-8") as f:
            f.write(cleaned_text)

        subprocess.run(['quarto', 'render', f"quarto_output/files/{fileName}.qmd" ], check=True)
        
        if os.path.exists(output_pptx_path):
            response =  FileResponse(
                open(output_pptx_path, 'rb'),
                as_attachment=True,
                filename=fileName,
            )
        else:
            raise Http404("Presentation file not found.")

        threading.Thread(target=cleanup, args=(input_qmd_path, output_pptx_path)).start()

        return response
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)
      
def download_from_media(request, filename):
     file_path = os.path.join(settings.MEDIA_ROOT, filename)
     try:
        return FileResponse(open(file_path,'rb'),as_attachment=True, filename=filename)
     except:
        raise Http404("File not found")
     
def cleanup(qmd_files, pptx_files):
    try:
        os.remove(qmd_files)
        os.remove(pptx_files)
    except:
        raise Http404("File not found")
